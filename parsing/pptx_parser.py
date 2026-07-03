""".pptx parser — single deterministic method: python-pptx."""
from pathlib import Path

from pptx import Presentation

from .common import ImageWriter, ParsedBlock, ParsedDocument

_PICTURE_SHAPE_TYPE = 13


def parse(path: Path, image_writer: ImageWriter) -> ParsedDocument:
    prs = Presentation(str(path))
    slide_count = len(prs.slides)
    parsed = ParsedDocument(source_file=path.name, doc_type="pptx", units_total=slide_count)
    parsed.metadata = {"slides": slide_count}

    for slide_idx, slide in enumerate(prs.slides, 1):
        title_shape = slide.shapes.title
        title_id = title_shape.shape_id if title_shape is not None else None
        if title_shape is not None and title_shape.text.strip():
            parsed.blocks.append(
                ParsedBlock(kind="heading", content=title_shape.text.strip(), level=2, unit=slide_idx)
            )
        for shape in slide.shapes:
            if shape.has_text_frame and shape.shape_id != title_id:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        parsed.blocks.append(ParsedBlock(kind="text", content=t, unit=slide_idx))
            if shape.shape_type == _PICTURE_SHAPE_TYPE:
                try:
                    blob = shape.image.blob
                    ext = shape.image.content_type.split("/")[-1].replace("jpeg", "jpg")
                    image_id = image_writer.add(blob, ext)
                    parsed.blocks.append(ParsedBlock(kind="image", image_id=image_id, unit=slide_idx))
                except Exception as e:
                    parsed.blocks.append(
                        ParsedBlock(kind="text", content=f"[image extraction error: {e}]", unit=slide_idx)
                    )
            if shape.has_table:
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                parsed.blocks.append(ParsedBlock(kind="table", content=rows, unit=slide_idx))

    return parsed
